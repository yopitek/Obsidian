#!/usr/bin/env python3
# /// script
# requires-python = ">=3.10"
# dependencies = [
#     "python-pptx>=1.0.0",
#     "Pillow>=10.0.0",
# ]
# ///
"""
Create PPT presentations from images.

Usage:
    uv run create_slides.py img1.png img2.png -o output.pptx
    uv run create_slides.py img1.png img2.png -t "Title 1" "Title 2" --layout title_above -o slides.pptx
    uv run create_slides.py *.png --layout grid --cols 3 -o gallery.pptx
"""

import argparse
import math
import sys
from pathlib import Path

from PIL import Image as PILImage
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt, Emu


LAYOUTS = ["fullscreen", "title_above", "title_below", "title_left", "center", "grid"]

SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)


def hex_to_rgb(hex_str: str) -> RGBColor:
    hex_str = hex_str.lstrip("#")
    r, g, b = int(hex_str[0:2], 16), int(hex_str[2:4], 16), int(hex_str[4:6], 16)
    return RGBColor(r, g, b)


def fit_image_size(
    img_w: int, img_h: int,
    max_w: int, max_h: int,
) -> tuple[int, int]:
    """Scale image to fit within max dimensions, preserving aspect ratio."""
    ratio_w = max_w / img_w
    ratio_h = max_h / img_h
    ratio = min(ratio_w, ratio_h)
    return int(img_w * ratio), int(img_h * ratio)


def cover_image_size(
    img_w: int, img_h: int,
    target_w: int, target_h: int,
) -> tuple[int, int, int, int]:
    """Scale image to cover target area (may crop). Returns (width, height, left_offset, top_offset)."""
    ratio_w = target_w / img_w
    ratio_h = target_h / img_h
    ratio = max(ratio_w, ratio_h)
    new_w = int(img_w * ratio)
    new_h = int(img_h * ratio)
    left = (target_w - new_w) // 2
    top = (target_h - new_h) // 2
    return new_w, new_h, left, top


def get_image_size(image_path: str) -> tuple[int, int]:
    with PILImage.open(image_path) as img:
        return img.size


def add_textbox(slide, text: str, left, top, width, height,
                font_size: int = 28, font_color: RGBColor = None,
                bold: bool = True, alignment=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.alignment = alignment
    if font_color:
        p.font.color.rgb = font_color
    return txBox


def set_slide_bg(slide, color: RGBColor):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_fullscreen_slide(prs, image_path: str, bg_color: RGBColor):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, bg_color)

    img_w, img_h = get_image_size(image_path)
    new_w, new_h, left, top = cover_image_size(
        img_w, img_h, SLIDE_WIDTH, SLIDE_HEIGHT
    )
    slide.shapes.add_picture(image_path, left, top, new_w, new_h)
    return slide


def add_title_image_slide(prs, image_path: str, title: str,
                          title_position: str, bg_color: RGBColor,
                          title_color: RGBColor, title_size: int, margin: float):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, bg_color)

    margin_emu = Inches(margin)
    title_h = Inches(1.2)
    content_w = SLIDE_WIDTH - 2 * margin_emu

    if title_position == "above":
        # Title on top
        if title:
            add_textbox(slide, title, margin_emu, margin_emu, content_w, title_h,
                        font_size=title_size, font_color=title_color, alignment=PP_ALIGN.LEFT)
        img_top = margin_emu + title_h + Inches(0.2)
        img_max_h = SLIDE_HEIGHT - img_top - margin_emu
    elif title_position == "below":
        # Title on bottom
        img_top = margin_emu
        img_max_h = SLIDE_HEIGHT - title_h - margin_emu * 2 - Inches(0.2)
        if title:
            add_textbox(slide, title, margin_emu, img_top + img_max_h + Inches(0.2),
                        content_w, title_h, font_size=title_size, font_color=title_color)
    else:  # left
        title_w = Inches(3.5)
        if title:
            add_textbox(slide, title, margin_emu, margin_emu, title_w, SLIDE_HEIGHT - 2 * margin_emu,
                        font_size=title_size, font_color=title_color)
        content_w = SLIDE_WIDTH - title_w - margin_emu * 3
        margin_emu = title_w + margin_emu * 2
        img_top = Inches(margin)
        img_max_h = SLIDE_HEIGHT - 2 * Inches(margin)

    # Add image
    img_w, img_h = get_image_size(image_path)
    new_w, new_h = fit_image_size(img_w, img_h, content_w, img_max_h)

    # Center horizontally within content area
    if title_position == "left":
        img_left = margin_emu + (content_w - new_w) // 2
    else:
        img_left = Inches(margin) + (content_w - new_w) // 2
    img_top_centered = img_top + (img_max_h - new_h) // 2

    slide.shapes.add_picture(image_path, img_left, img_top_centered, new_w, new_h)
    return slide


def add_center_slide(prs, image_path: str, bg_color: RGBColor, margin: float):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, bg_color)

    margin_emu = Inches(margin)
    max_w = SLIDE_WIDTH - 2 * margin_emu
    max_h = SLIDE_HEIGHT - 2 * margin_emu

    img_w, img_h = get_image_size(image_path)
    new_w, new_h = fit_image_size(img_w, img_h, max_w, max_h)

    left = (SLIDE_WIDTH - new_w) // 2
    top = (SLIDE_HEIGHT - new_h) // 2
    slide.shapes.add_picture(image_path, left, top, new_w, new_h)
    return slide


def add_grid_slide(prs, image_paths: list[str], cols: int,
                   bg_color: RGBColor, margin: float):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, bg_color)

    margin_emu = Inches(margin)
    gap = Inches(0.2)
    rows = math.ceil(len(image_paths) / cols)

    total_w = SLIDE_WIDTH - 2 * margin_emu - (cols - 1) * gap
    total_h = SLIDE_HEIGHT - 2 * margin_emu - (rows - 1) * gap
    cell_w = total_w // cols
    cell_h = total_h // rows

    for idx, img_path in enumerate(image_paths):
        row = idx // cols
        col = idx % cols

        cell_left = margin_emu + col * (cell_w + gap)
        cell_top = margin_emu + row * (cell_h + gap)

        img_w, img_h = get_image_size(img_path)
        new_w, new_h = fit_image_size(img_w, img_h, cell_w, cell_h)

        img_left = cell_left + (cell_w - new_w) // 2
        img_top = cell_top + (cell_h - new_h) // 2
        slide.shapes.add_picture(img_path, img_left, img_top, new_w, new_h)

    return slide


def main():
    parser = argparse.ArgumentParser(description="Create PPT from images")
    parser.add_argument("images", nargs="+", help="Image file paths")
    parser.add_argument("-o", "--output", required=True, help="Output .pptx file path")
    parser.add_argument("-t", "--titles", nargs="*", help="Slide titles")
    parser.add_argument("--layout", choices=LAYOUTS, default="fullscreen",
                        help="Layout mode (default: fullscreen)")
    parser.add_argument("--bg-color", default="FFFFFF", help="Background color hex (default: FFFFFF)")
    parser.add_argument("--title-color", default="333333", help="Title font color hex (default: 333333)")
    parser.add_argument("--title-size", type=int, default=28, help="Title font size in pt")
    parser.add_argument("--cols", type=int, default=2, help="Grid columns (for grid layout)")
    parser.add_argument("--margin", type=float, default=0.5, help="Margin in inches")
    parser.add_argument("--template", help="Template .pptx file path")

    args = parser.parse_args()

    # Validate images
    image_paths = []
    for img in args.images:
        p = Path(img)
        if not p.exists():
            print(f"Warning: Image not found, skipping: {img}", file=sys.stderr)
            continue
        image_paths.append(str(p))

    if not image_paths:
        print("Error: No valid image files found", file=sys.stderr)
        sys.exit(1)

    # Create presentation
    if args.template:
        prs = Presentation(args.template)
    else:
        prs = Presentation()
        prs.slide_width = SLIDE_WIDTH
        prs.slide_height = SLIDE_HEIGHT

    bg_color = hex_to_rgb(args.bg_color)
    title_color = hex_to_rgb(args.title_color)
    titles = args.titles or []

    if args.layout == "grid":
        # Grid: batch images into slides
        for i in range(0, len(image_paths), args.cols * 2):
            batch = image_paths[i:i + args.cols * 2]
            add_grid_slide(prs, batch, args.cols, bg_color, args.margin)
    else:
        for idx, img_path in enumerate(image_paths):
            title = titles[idx] if idx < len(titles) else None

            if args.layout == "fullscreen":
                add_fullscreen_slide(prs, img_path, bg_color)
            elif args.layout == "center":
                add_center_slide(prs, img_path, bg_color, args.margin)
            elif args.layout in ("title_above", "title_below", "title_left"):
                position = args.layout.replace("title_", "")
                add_title_image_slide(
                    prs, img_path, title, position,
                    bg_color, title_color, args.title_size, args.margin
                )

    # Save
    output_path = Path(args.output)
    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output_path))
    print(f"Presentation saved: {output_path.resolve()}")
    print(f"  {len(image_paths)} images, {len(prs.slides)} slides, layout: {args.layout}")


if __name__ == "__main__":
    main()

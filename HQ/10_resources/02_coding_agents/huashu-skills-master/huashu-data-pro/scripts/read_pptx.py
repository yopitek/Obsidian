#!/usr/bin/env python3
"""
read_pptx.py — 读取 PPTX 文件并输出结构化内容

用法：
    python3 read_pptx.py presentation.pptx
    python3 read_pptx.py presentation.pptx --format markdown
    python3 read_pptx.py presentation.pptx --thumbnails ./thumbs/
    python3 read_pptx.py presentation.pptx --inventory

参数：
    文件路径              PPTX 文件
    --format FMT         输出格式：text（默认）/ markdown / json
    --thumbnails DIR     将每页导出为缩略图到指定目录
    --inventory          仅输出结构信息（页数、布局、尺寸）
    --slide N            仅输出第 N 页（从1开始）

依赖：
    pip install python-pptx Pillow
"""

import sys
import json
import argparse
import os


def ensure_dependencies():
    """检查并提示安装依赖"""
    missing = []
    try:
        import pptx
    except ImportError:
        missing.append('python-pptx')
    if missing:
        print(f"缺少依赖: {', '.join(missing)}", file=sys.stderr)
        print(f"请运行: pip install {' '.join(missing)}", file=sys.stderr)
        sys.exit(1)

ensure_dependencies()

from pptx import Presentation
from pptx.util import Inches, Pt, Emu


def read_pptx(filepath, output_format='text', inventory_only=False, slide_num=None):
    """读取 PPTX 并格式化输出"""
    try:
        prs = Presentation(filepath)
    except Exception as e:
        print(f"读取失败: {e}", file=sys.stderr)
        sys.exit(1)

    width_in = prs.slide_width / 914400
    height_in = prs.slide_height / 914400
    total_slides = len(prs.slides)

    print(f"文件: {os.path.basename(filepath)}")
    print(f"尺寸: {width_in:.1f}\" x {height_in:.1f}\" ({width_in/height_in:.2f}:1)")
    print(f"页数: {total_slides}")

    if inventory_only:
        print(f"\n幻灯片清单:")
        for i, slide in enumerate(prs.slides, 1):
            layout_name = slide.slide_layout.name if slide.slide_layout else "无布局"
            shapes_count = len(slide.shapes)
            has_notes = bool(slide.notes_slide and slide.notes_slide.notes_text_frame.text.strip())
            text_preview = ""
            for shape in slide.shapes:
                if shape.has_text_frame:
                    text_preview = shape.text_frame.text.strip()[:60]
                    if text_preview:
                        break
            print(f"  [{i:>3}] 布局: {layout_name:<20} 元素: {shapes_count:<4} 备注: {'有' if has_notes else '无'}  {text_preview}")
        return

    # 提取内容
    slides_data = []
    for i, slide in enumerate(prs.slides, 1):
        if slide_num and i != slide_num:
            continue

        slide_data = {
            'number': i,
            'layout': slide.slide_layout.name if slide.slide_layout else None,
            'elements': []
        }

        for shape in slide.shapes:
            element = {
                'type': shape.shape_type,
                'name': shape.name,
                'position': {
                    'left': round(shape.left / 914400, 2) if shape.left else 0,
                    'top': round(shape.top / 914400, 2) if shape.top else 0,
                    'width': round(shape.width / 914400, 2) if shape.width else 0,
                    'height': round(shape.height / 914400, 2) if shape.height else 0
                }
            }

            if shape.has_text_frame:
                paragraphs = []
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if text:
                        paragraphs.append(text)
                element['text'] = paragraphs

            if shape.has_table:
                table = shape.table
                rows = []
                for row in table.rows:
                    cells = [cell.text.strip() for cell in row.cells]
                    rows.append(cells)
                element['table'] = rows

            if hasattr(shape, 'image'):
                try:
                    element['image'] = {
                        'content_type': shape.image.content_type,
                        'size': len(shape.image.blob)
                    }
                except Exception:
                    pass

            slide_data['elements'].append(element)

        # 提取备注
        if slide.notes_slide and slide.notes_slide.notes_text_frame.text.strip():
            slide_data['notes'] = slide.notes_slide.notes_text_frame.text.strip()

        slides_data.append(slide_data)

    # 输出
    if output_format == 'json':
        print(json.dumps(slides_data, ensure_ascii=False, indent=2))
    elif output_format == 'markdown':
        for slide_data in slides_data:
            print(f"\n## 第 {slide_data['number']} 页")
            if slide_data.get('layout'):
                print(f"*布局: {slide_data['layout']}*\n")
            for el in slide_data['elements']:
                if 'text' in el and el['text']:
                    for para in el['text']:
                        print(f"- {para}")
                if 'table' in el:
                    table = el['table']
                    if len(table) > 0:
                        # 表头
                        print(f"\n| {'|'.join(table[0])} |")
                        print(f"|{'|'.join(['---'] * len(table[0]))}|")
                        for row in table[1:]:
                            print(f"| {'|'.join(row)} |")
                        print()
            if slide_data.get('notes'):
                print(f"\n> 备注: {slide_data['notes']}")
    else:
        for slide_data in slides_data:
            print(f"\n{'='*50}")
            print(f"第 {slide_data['number']} 页 (布局: {slide_data.get('layout', '无')})")
            print(f"{'='*50}")
            for el in slide_data['elements']:
                if 'text' in el and el['text']:
                    for para in el['text']:
                        print(f"  {para}")
                if 'table' in el:
                    for row in el['table']:
                        print(f"  | {' | '.join(row)} |")
            if slide_data.get('notes'):
                print(f"  [备注] {slide_data['notes']}")


def export_thumbnails(filepath, output_dir):
    """导出缩略图（需要 Pillow）"""
    try:
        from PIL import Image
        import io
    except ImportError:
        print("导出缩略图需要 Pillow: pip install Pillow", file=sys.stderr)
        sys.exit(1)

    # 缩略图方案：用 Playwright 截图每页 HTML（如果有源文件）
    # 否则提取 PPT 中的图片作为参考
    print("提示：PPTX 原生缩略图导出建议使用 LibreOffice：")
    print(f"  libreoffice --headless --convert-to png {filepath}")
    print(f"或使用 Playwright 截图 HTML 源文件（如果有的话）")


def main():
    parser = argparse.ArgumentParser(description='读取 PPTX 文件')
    parser.add_argument('filepath', help='PPTX 文件路径')
    parser.add_argument('--format', choices=['text', 'markdown', 'json'], default='text', help='输出格式')
    parser.add_argument('--inventory', action='store_true', help='仅输出结构信息')
    parser.add_argument('--thumbnails', help='导出缩略图到指定目录')
    parser.add_argument('--slide', type=int, help='仅输出第 N 页')

    args = parser.parse_args()

    if args.thumbnails:
        export_thumbnails(args.filepath, args.thumbnails)
    else:
        read_pptx(args.filepath, args.format, args.inventory, args.slide)


if __name__ == '__main__':
    main()
